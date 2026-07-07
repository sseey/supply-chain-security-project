#!/usr/bin/env python3
"""
Génère le PowerPoint de COURS (30 min) sur la sécurité de la chaîne
d'approvisionnement logicielle, avec présentation des outils ET des
NOTES DU PRÉSENTATEUR sur chaque slide.

À distinguer du deck d'intro projet (generate_slides.py) : ici c'est un cours
magistral court, orienté concepts + démonstration d'outils.

Usage :
    pip install python-pptx
    python generate_cours_slides.py
    # -> cours-supply-chain-security.pptx  (ignoré par git)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ------------------------------------------------------------------- palette (idem intro)
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
NAVY2  = RGBColor(0x13, 0x2A, 0x4D)
TEAL   = RGBColor(0x10, 0xB9, 0x81)
TEALD  = RGBColor(0x0B, 0x81, 0x5A)
RED    = RGBColor(0xEF, 0x44, 0x44)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
INK    = RGBColor(0x1F, 0x29, 0x37)
GREY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG = RGBColor(0x0B, 0x1F, 0x3A)
CODEFG = RGBColor(0xC9, 0xF3, 0xE4)
FONT   = "Calibri"
MONO   = "Consolas"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=5):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz)
            r.font.color.rgb = col; r.font.bold = bold; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=18, color=INK, gap=9):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        lvl, txt = (it if isinstance(it, tuple) else (0, it))
        r = p.add_run(); r.text = ("▸  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size if lvl == 0 else size - 3)
        r.font.color.rgb = color if lvl == 0 else GREY
        r.font.name = FONT
    return tb


def header(s, kicker, title, accent=TEAL):
    bg(s, WHITE)
    rect(s, 0, 0, Inches(0.28), H, accent)
    text(s, Inches(0.6), Inches(0.42), Inches(12), Inches(0.4), [(kicker.upper(), 13, accent, True)])
    text(s, Inches(0.6), Inches(0.74), Inches(12.2), Inches(0.9), [(title, 29, NAVY, True)])
    rect(s, Inches(0.62), Inches(1.55), Inches(1.7), Pt(3), accent)


def footer(s, n, dark=False):
    col = RGBColor(0xB8, 0xC2, 0xD0) if dark else GREY
    text(s, Inches(0.6), Inches(7.03), Inches(9), Inches(0.35),
         [("Cours — Sécurité de la chaîne d'approvisionnement logicielle (SLSA)", 9, col, False)])
    text(s, Inches(11.6), Inches(7.03), Inches(1.4), Inches(0.35), [(str(n), 9, col, False)], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, label, color, fs=12):
    c = rect(s, x, y, w, Inches(0.42), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = c.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(fs); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT
    return c


def code(s, x, y, w, h, lines):
    rect(s, x, y, w, h, CODEBG, MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), h - Inches(0.24))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        comment = ln.strip().startswith("#")
        r = p.add_run(); r.text = ln
        r.font.name = MONO; r.font.size = Pt(12.5)
        r.font.color.rgb = RGBColor(0x8A, 0xA0, 0xB4) if comment else CODEFG
    return tb


def toolcard(s, x, y, w, h, name, role, color):
    rect(s, x, y, w, h, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, w, Inches(0.5), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.2), y + Inches(0.07), w - Inches(0.4), Inches(0.4), [(name, 15, WHITE, True)])
    text(s, x + Inches(0.2), y + Inches(0.6), w - Inches(0.4), h - Inches(0.7), [(role, 12.5, INK, False)])


def notes(s, lines):
    """Ajoute les notes du présentateur (lignes = liste de str)."""
    tf = s.notes_slide.notes_text_frame
    tf.text = lines[0]
    for ln in lines[1:]:
        p = tf.add_paragraph(); p.text = ln


# =================================================================== 1. TITRE
s = slide(); bg(s, NAVY)
rect(s, 0, 0, W, Inches(0.16), TEAL); rect(s, 0, Inches(7.34), W, Inches(0.16), TEAL)
text(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.5), [("COURS · 30 MINUTES", 15, TEAL, True)])
text(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(2.0),
     [[("Sécuriser la chaîne", 46, WHITE, True)],
      [("d'approvisionnement logicielle", 46, WHITE, True)]])
text(s, Inches(0.9), Inches(4.3), Inches(11.4), Inches(0.8),
     [("Concepts, référentiel SLSA, et démonstration des outils : Syft · Grype · cosign · Kyverno", 18, RGBColor(0xC7,0xD2,0xE0), False)])
chip(s, Inches(0.9), Inches(5.5), Inches(2.4), "SBOM", TEALD)
chip(s, Inches(3.5), Inches(5.5), Inches(2.4), "Signature", TEALD)
chip(s, Inches(6.1), Inches(5.5), Inches(2.4), "Provenance", TEALD)
chip(s, Inches(8.7), Inches(5.5), Inches(2.9), "Admission control", TEALD)
notes(s, [
    "⏱️ 0:00–1:00 (1 min).",
    "Accroche : « Vous savez tous faire un pipeline CI/CD. Aujourd'hui on répond à LA question "
    "que se posent les équipes DevSecOps : comment prouver que l'image en prod est bien la nôtre, "
    "non altérée ? »",
    "Annoncer le format : 30 min, orienté concepts + démo d'outils. On finira par une chaîne complète.",
    "Poser une question à la salle : « qui a déjà signé une image conteneur ? » pour jauger le niveau.",
])

# =================================================================== 2. AGENDA
s = slide(); header(s, "Plan du cours", "Ce qu'on va voir en 30 minutes")
items = [
    ("1", "Le problème & les attaques réelles", "pourquoi la supply chain est la nouvelle cible", TEAL, "~5 min"),
    ("2", "SBOM & scan — Syft, Grype", "connaître et auditer ce qu'on livre", TEALD, "~6 min"),
    ("3", "Signature — cosign / Sigstore", "prouver l'origine (keyless, Rekor)", NAVY2, "~7 min"),
    ("4", "Attestations & SLSA", "provenance vérifiable, niveaux de maturité", AMBER, "~5 min"),
    ("5", "Admission control — Kyverno", "le cluster qui refuse l'inconnu", RED, "~5 min"),
    ("6", "Chaîne complète & panorama outils", "récap + alternatives", GREY, "~2 min"),
]
y = Inches(1.9)
for num, t, d, col, dur in items:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(0.78), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    b = rect(s, Inches(0.72), y + Inches(0.13), Inches(0.52), Inches(0.52), col, MSO_SHAPE.OVAL)
    tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = WHITE; r.font.name = FONT
    text(s, Inches(1.5), y + Inches(0.08), Inches(6.8), Inches(0.5), [(t, 16, NAVY, True)])
    text(s, Inches(1.5), y + Inches(0.44), Inches(7.5), Inches(0.35), [(d, 12, GREY, False)])
    text(s, Inches(11.2), y + Inches(0.2), Inches(1.4), Inches(0.4), [(dur, 13, col, True)], align=PP_ALIGN.RIGHT)
    y = y + Inches(0.86)
notes(s, [
    "⏱️ 1:00–2:00 (1 min).",
    "Donner la carte du cours pour que la salle situe chaque partie. Insister : on va MANIPULER des "
    "outils, pas seulement en parler.",
    "Fil conducteur à répéter tout le long : on passe de « on scanne et on espère » à « on vérifie et on bloque ».",
])

# =================================================================== 3. LE PROBLÈME
s = slide(); header(s, "Partie 1 — le problème", "La supply chain logicielle : nouvelle surface d'attaque", RED)
text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.6),
     [[("La chaîne d'appro = ", 17, INK, False), ("tout ce qui transforme le code source en artefact qui tourne", 17, NAVY, True),
       (" : dépendances, build, registry, déploiement.", 17, INK, False)]])
cards = [
    ("SolarWinds — 2020", "Malware injecté dans le BUILD, signé, distribué à 18 000 clients."),
    ("Codecov — 2021", "Script CI modifié → vol des secrets de milliers de pipelines."),
    ("Dependency confusion", "Faux paquets internes publiés sur les registries publics."),
    ("XZ Utils — 2024", "Backdoor glissée sur 3 ans dans une dépendance open source."),
]
x0, y0, cw, ch, gx, gy = Inches(0.6), Inches(2.6), Inches(5.95), Inches(1.5), Inches(0.2), Inches(0.2)
for i, (t, d) in enumerate(cards):
    cx = x0 + (cw + gx) * (i % 2); cy = y0 + (ch + gy) * (i // 2)
    rect(s, cx, cy, cw, ch, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE); rect(s, cx, cy, Inches(0.12), ch, RED)
    text(s, cx + Inches(0.32), cy + Inches(0.14), cw - Inches(0.5), Inches(0.5), [(t, 16, NAVY, True)])
    text(s, cx + Inches(0.32), cy + Inches(0.6), cw - Inches(0.5), Inches(0.8), [(d, 13, INK, False)])
text(s, Inches(0.6), Inches(6.05), Inches(12), Inches(0.6),
     [[("Point clé : ", 16, RED, True), ("un « docker pull » ne vérifie rien, et un scan vert ne prouve pas que l'image DÉPLOYÉE = image scannée.", 16, INK, False)]])
footer(s, 3)
notes(s, [
    "⏱️ 2:00–5:00 (3 min).",
    "Raconter 1 ou 2 attaques en histoire, pas en liste. SolarWinds : l'attaquant n'a pas touché au code "
    "source visible, il a compromis le SERVEUR DE BUILD ; l'artefact malveillant était signé par l'éditeur → "
    "confiance trahie en aval sur 18 000 orgs.",
    "XZ (2024) : montrer que même l'open source très surveillé est vulnérable ; backdoor introduite par un "
    "mainteneur patient sur 3 ans. Découverte par hasard (perf ssh).",
    "Conclure sur le point clé du bas : nos réflexes habituels (scan, pull) ne suffisent pas. "
    "Transition : « il faut rendre chaque maillon PROUVABLE ». → partie 2.",
])

# =================================================================== 4. MODÈLE MENTAL
s = slide(); header(s, "Le changement de posture", "De « on scanne » à « on vérifie et on bloque »")
# gauche : avant
rect(s, Inches(0.6), Inches(1.9), Inches(5.9), Inches(4.4), RGBColor(0xFB,0xE9,0xE9), MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.9), Inches(2.1), Inches(5.3), Inches(0.5), [("❌  Approche classique", 18, RED, True)])
bullets(s, Inches(0.95), Inches(2.75), Inches(5.3), Inches(3.3), [
    "On scanne l'image dans la CI", "« C'est vert » → on déploie",
    "Confiance implicite au registry et au tag", "Aucune preuve à l'exécution",
    "L'image déployée PEUT différer de la scannée",
], size=15, color=INK, gap=13)
# droite : après
rect(s, Inches(6.8), Inches(1.9), Inches(5.9), Inches(4.4), RGBColor(0xE7,0xF7,0xF0), MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(0.5), [("✅  Chaîne vérifiable (zero-trust)", 18, TEALD, True)])
bullets(s, Inches(7.15), Inches(2.75), Inches(5.3), Inches(3.3), [
    "On produit des PREUVES (SBOM, signature, provenance)",
    "Preuves liées au DIGEST, pas au tag",
    "Le cluster VÉRIFIE avant d'exécuter",
    "Ce qui n'est pas prouvé est REFUSÉ",
    "Chaque garantie est reproductible par une commande",
], size=15, color=INK, gap=13)
footer(s, 4)
notes(s, [
    "⏱️ 5:00–7:00 (2 min).",
    "C'est LE slide conceptuel. Marteler : la sécurité n'est plus une étape (« le scan »), c'est une "
    "PROPRIÉTÉ vérifiable de bout en bout.",
    "Insister sur digest vs tag : un tag est mutable (on peut repointer :v1 vers une autre image) ; le "
    "digest sha256 est le hash du contenu, immuable. Toutes nos preuves s'accrochent au digest.",
    "Annoncer les 4 briques qui suivent : SBOM, signature, attestations, admission.",
])

# =================================================================== 5. SBOM / SYFT
s = slide(); header(s, "Partie 2 — outil", "SBOM avec Syft : l'inventaire de ce qu'on livre", TEALD)
text(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.6),
     [[("SBOM", 16, TEALD, True), (" (Software Bill of Materials) = la liste exhaustive des paquets/versions d'une image. Son ", 15, INK, False),
       ("« étiquette de composition »", 15, NAVY, True), (".", 15, INK, False)]])
toolcard(s, Inches(0.6), Inches(2.55), Inches(3.5), Inches(2.0), "Syft", "Génère le SBOM d'une image ou d'un filesystem. Formats standards SPDX et CycloneDX.", TEALD)
code(s, Inches(4.35), Inches(2.55), Inches(8.35), Inches(2.0), [
    "# Générer le SBOM (format SPDX)",
    "syft ghcr.io/moi/app:1.0 -o spdx-json > sbom.spdx.json",
    "",
    "# Aperçu lisible",
    "syft ghcr.io/moi/app:1.0 -o table",
])
text(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.5),
     [[("Pourquoi c'est utile : ", 15, NAVY, True),
       ("le jour où une CVE tombe sur une lib, « suis-je affecté ? » devient répondable en SECONDES "
        "en cherchant dans les SBOM. C'est aussi une exigence réglementaire montante (US EO 14028, Cyber Resilience Act).", 15, INK, False)]])
footer(s, 5)
notes(s, [
    "⏱️ 7:00–9:30 (2,5 min).",
    "Démo live possible : lancer `syft` sur l'image du projet et montrer Python/Flask/gunicorn dans la liste.",
    "Analogie : le SBOM est la liste des ingrédients sur un emballage alimentaire. Sans elle, en cas de "
    "rappel d'un ingrédient, impossible de savoir quels produits sont concernés.",
    "Mentionner les 2 formats (SPDX = ISO, orienté conformité ; CycloneDX = OWASP, orienté sécurité). "
    "Syft génère aussi le SBOM DANS le build (bonne pratique : SBOM au moment du build, pas après).",
])

# =================================================================== 6. GRYPE
s = slide(); header(s, "Partie 2 — outil", "Scan de vulnérabilités avec Grype (et la gate)", TEALD)
toolcard(s, Inches(0.6), Inches(1.85), Inches(3.5), Inches(2.0), "Grype", "Scanne une image ou un SBOM contre les bases de CVE. Peut CASSER le build selon un seuil.", AMBER)
code(s, Inches(4.35), Inches(1.85), Inches(8.35), Inches(2.0), [
    "# Scanner le SBOM déjà produit",
    "grype sbom:sbom.spdx.json",
    "",
    "# Gate CI : échoue (exit != 0) sur CVE critique corrigeable",
    "grype ghcr.io/moi/app:1.0 --only-fixed --fail-on critical",
])
text(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.55),
     [[("L'idée clé : ", 15, NAVY, True), ("le scan ne se contente pas d'afficher — il ", 15, INK, False),
       ("ARRÊTE la chaîne", 15, RED, True), (" quand une vulnérabilité actionnable est présente.", 15, INK, False)]])
chip(s, Inches(0.6), Inches(4.85), Inches(3.9), "Alternative : Trivy (Aqua)", NAVY2, fs=12)
chip(s, Inches(4.7), Inches(4.85), Inches(4.0), "only-fixed = moins de bruit", NAVY2, fs=12)
text(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.9),
     [[("Bonne pratique : ", 15, TEALD, True), ("ne bloquer que sur le corrigeable (", 15, INK, False),
       ("--only-fixed", 13, TEALD, True), ("), sinon les équipes désactivent la gate à force de faux positifs non actionnables.", 15, INK, False)]])
footer(s, 6)
notes(s, [
    "⏱️ 9:30–11:30 (2 min).",
    "Démo : introduire volontairement une vieille version de Flask, rebuild, relancer grype → montrer la "
    "gate qui casse (exit code != 0). C'est très parlant.",
    "Débat à lancer : faut-il bloquer sur TOUTES les CVE ? Non → sinon bruit ingérable, on bloque sur "
    "critical/high CORRIGEABLE. Une CVE sans correctif dispo : la bloquer n'aide pas, il faut une exception tracée.",
    "Comparer vite Grype et Trivy : mêmes usages, choisir selon l'écosystème. Ne pas s'attarder.",
])

# =================================================================== 7. SIGSTORE / COSIGN
s = slide(); header(s, "Partie 3 — outil", "Signer l'image : cosign & Sigstore", NAVY2)
text(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.5),
     [[("Signer = attacher une preuve cryptographique « c'est bien NOUS qui l'avons produite ». ", 15, INK, False)]])
# 3 composants Sigstore
comps = [("cosign", "signe/vérifie images & attestations", TEALD),
         ("Fulcio", "délivre un certificat éphémère lié à votre identité OIDC", NAVY2),
         ("Rekor", "journal de transparence PUBLIC et immuable des signatures", AMBER)]
x = Inches(0.6)
for name, role, col in comps:
    toolcard(s, x, Inches(2.35), Inches(3.95), Inches(1.55), name, role, col)
    x = x + Inches(4.05)
code(s, Inches(0.6), Inches(4.15), Inches(12.1), Inches(1.55), [
    "# Signature KEYLESS : pas de clé à gérer, identité via OIDC (GitHub/Google)",
    "cosign sign ghcr.io/moi/app@sha256:...          # signe PAR DIGEST",
    "cosign verify --certificate-identity ... --certificate-oidc-issuer ...  ghcr.io/moi/app@sha256:...",
])
text(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.7),
     [[("Keyless > par clé : ", 15, TEALD, True),
       ("rien à stocker ni faire fuiter ; l'identité (workflow CI, branche) est vérifiable et la preuve est publique dans Rekor.", 15, INK, False)]])
footer(s, 7)
notes(s, [
    "⏱️ 11:30–15:30 (4 min) — le cœur du cours, prendre le temps.",
    "Expliquer le problème des clés : une signature par clé privée classique suppose de gérer/stocker/faire "
    "tourner la clé → surface de fuite. Sigstore supprime la clé longue durée.",
    "Keyless en 3 temps : (1) cosign obtient un jeton OIDC prouvant votre identité ; (2) Fulcio émet un "
    "certificat valable ~10 min lié à cette identité ; (3) la signature + le certificat sont inscrits dans "
    "Rekor (log public immuable, type blockchain/Merkle). N'importe qui peut auditer.",
    "En CI, l'identité devient « ce workflow, sur cette branche » → on peut exiger EXACTEMENT ça à la vérif.",
    "Démo si possible : cosign sign keyless (ouvre navigateur OIDC), puis cosign tree pour voir la signature attachée.",
    "Insister : on signe TOUJOURS par digest, jamais par tag mutable.",
])

# =================================================================== 8. ATTESTATIONS / PROVENANCE
s = slide(); header(s, "Partie 4 — concept", "Attestations & provenance : prouver l'ORIGINE", AMBER)
text(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.6),
     [[("Attestation = ", 16, AMBER, True), ("une affirmation SIGNÉE attachée à l'image. ", 15, INK, False),
       ("Deux exemples clés : le SBOM lui-même, et la PROVENANCE.", 15, NAVY, True)]])
rect(s, Inches(0.6), Inches(2.55), Inches(5.9), Inches(1.9), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.85), Inches(2.7), Inches(5.4), Inches(1.7),
     [[("Provenance ", 15, NAVY, True), ("répond à :", 15, INK, False)],
      [("qui a construit ?", 14, INK, False)], [("depuis quel commit / source ?", 14, INK, False)],
      [("avec quel builder ?  quand ?", 14, INK, False)]], sa=4)
code(s, Inches(6.7), Inches(2.55), Inches(6.0), Inches(1.9), [
    "# Attacher le SBOM et la provenance",
    "cosign attest --type spdxjson \\",
    "   --predicate sbom.spdx.json  IMG@sha256:...",
    "cosign attest --type slsaprovenance \\",
    "   --predicate provenance.json IMG@sha256:...",
])
text(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(0.5),
     [[("Ces attestations sont ", 15, INK, False), ("stockées à côté de l'image dans le registry", 15, NAVY, True),
       (" et vérifiables (cosign verify-attestation).", 15, INK, False)]])
chip(s, Inches(0.6), Inches(5.35), Inches(5.6), "En CI : provenance générée automatiquement", TEALD, fs=12)
text(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.6),
     [[("C'est cette provenance signée, produite par une plateforme de build, qui fait grimper le niveau ", 14, INK, False),
       ("SLSA", 14, TEALD, True), (" (slide suivant).", 14, INK, False)]])
footer(s, 8)
notes(s, [
    "⏱️ 15:30–18:30 (3 min).",
    "Distinguer signature (« c'est nous ») et attestation (« et voici des FAITS signés sur cette image »). "
    "La signature dit qui ; l'attestation dit quoi/comment.",
    "La provenance est le document le plus important : c'est le « certificat de naissance » de l'image. "
    "En cas d'incident, elle permet de remonter au commit et au build exacts.",
    "Souligner : à la main c'est pédagogique, mais la VRAIE valeur vient d'une provenance générée par la "
    "plateforme de CI (non falsifiable par le dev) → transition SLSA.",
])

# =================================================================== 9. SLSA
s = slide(); header(s, "Partie 4 — référentiel", "SLSA : des niveaux de garantie sur la provenance", TEALD)
text(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.5),
     [[("SLSA", 16, TEALD, True), (" — Supply-chain Levels for Software Artifacts (OpenSSF). Un référentiel de MATURITÉ.", 15, INK, False)]])
levels = [
    ("L1", "La provenance existe", "le build enregistre comment l'artefact a été fait", RGBColor(0xCF,0xE9,0xDD)),
    ("L2", "Provenance signée + build hébergé", "build sur une plateforme (pas un poste), provenance signée", TEAL),
    ("L3", "Build renforcé & isolé", "isolation forte, provenance infalsifiable et non contournable", RGBColor(0xCF,0xE9,0xDD)),
]
y = Inches(2.45)
for lvl, t, d, col in levels:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(1.15), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    b = rect(s, Inches(0.8), y + Inches(0.2), Inches(1.15), Inches(0.75), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lvl; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = FONT
    text(s, Inches(2.2), y + Inches(0.16), Inches(10.3), Inches(0.5), [(t, 17, NAVY, True)])
    text(s, Inches(2.2), y + Inches(0.62), Inches(10.3), Inches(0.5), [(d, 13.5, INK, False)])
    y = y + Inches(1.28)
text(s, Inches(2.2), Inches(3.55), Inches(10), Inches(0.4), [[("← cible réaliste d'un pipeline GitHub Actions", 12.5, TEALD, True)]])
footer(s, 9)
notes(s, [
    "⏱️ 18:30–20:30 (2 min).",
    "SLSA se prononce « salsa ». Ce n'est PAS un outil, c'est un cadre pour se situer et progresser.",
    "L1 = tu as une provenance. L2 = elle est signée ET le build tourne sur une plateforme (pas sur le "
    "laptop d'un dev, condition majeure). L3 = le build est isolé, la provenance non falsifiable même par "
    "un mainteneur.",
    "Message honnête : avec GitHub Actions + cosign keyless on vise L2. L3 demande des générateurs isolés "
    "dédiés (ex. slsa-github-generator). Demander aux étudiants d'annoncer leur niveau RÉEL dans le rapport.",
])

# =================================================================== 10. KYVERNO
s = slide(); header(s, "Partie 5 — outil", "Admission control : Kyverno, le gardien du cluster", RED)
text(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.6),
     [[("Un ", 15, INK, False), ("admission controller", 15, NAVY, True),
       (" intervient AVANT la création d'un Pod : il peut valider, muter… ou REFUSER la requête.", 15, INK, False)]])
toolcard(s, Inches(0.6), Inches(2.45), Inches(3.5), Inches(2.1), "Kyverno", "Politiques Kubernetes en YAML. verifyImages vérifie signatures & attestations et bloque le non conforme.", RED)
code(s, Inches(4.35), Inches(2.45), Inches(8.35), Inches(2.1), [
    "verifyImages:",
    "  - imageReferences: ['ghcr.io/moi/app*']",
    "    attestors:",
    "      - entries: [{ keyless: { issuer: ..., subject: ... }}]",
    "# validationFailureAction: Enforce  <-- REFUSE (vs Audit = journalise)",
])
text(s, Inches(0.6), Inches(4.75), Inches(12.1), Inches(0.55),
     [[("Le réglage qui change tout : ", 15, RED, True),
       ("Enforce", 15, RED, True), (" bloque réellement ; ", 15, INK, False),
       ("Audit", 15, AMBER, True), (" ne fait que journaliser (utile pour un déploiement progressif).", 15, INK, False)]])
chip(s, Inches(0.6), Inches(5.45), Inches(5.9), "Alternatives : OPA/Gatekeeper, Sigstore policy-controller", NAVY2, fs=12)
text(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.6),
     [[("On exige aussi : registry autorisé, déploiement par digest, présence d'attestation de provenance.", 14, INK, False)]])
footer(s, 10)
notes(s, [
    "⏱️ 20:30–24:00 (3,5 min).",
    "C'est le maillon qui rend tout le reste UTILE : sans vérification à l'admission, produire des signatures "
    "ne protège rien à l'exécution.",
    "Expliquer le webhook : l'API server appelle Kyverno avant de persister le Pod ; Kyverno va chercher la "
    "signature/attestation dans le registry et vérifie contre l'identité attendue.",
    "Insister sur Enforce vs Audit : erreur classique = croire qu'on est protégé alors que la politique est "
    "en Audit. Toujours vérifier. Conseil terrain : démarrer en Audit, observer, puis passer Enforce.",
    "Démo forte : `kubectl run` d'une image non signée → refus immédiat avec message Kyverno. C'est le "
    "moment « waouh » du cours.",
])

# =================================================================== 11. CHAÎNE COMPLÈTE
s = slide(); header(s, "Partie 6 — synthèse", "La chaîne complète, de bout en bout")
steps = [("Build", TEAL), ("SBOM", TEAL), ("Scan", AMBER), ("Sign", TEALD), ("Attest", TEALD), ("Push", NAVY2)]
x = Inches(0.6); y = Inches(2.0); cw = Inches(1.72)
for i, (lbl, col) in enumerate(steps):
    chip(s, x, y, cw, lbl, col, fs=13)
    x = x + cw + Inches(0.05)
    if i < len(steps) - 1:
        a = rect(s, x, y + Inches(0.02), Inches(0.3), Inches(0.38), GREY, MSO_SHAPE.CHEVRON); x = x + Inches(0.35)
rect(s, Inches(0.6), Inches(2.75), Inches(12.1), Inches(2.0), NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.5),
     [[("À l'admission, Kyverno vérifie — sinon ", 17, WHITE, True), ("❌ REFUSÉ", 17, RED, True), (" :", 17, WHITE, True)]])
checks = ["Signée par notre identité", "Provenance présente", "Registry autorisé + digest", "Pas de :latest / CVE critique"]
for i, c in enumerate(checks):
    bx = Inches(0.9) + Inches(2.95) * i
    rect(s, bx, Inches(3.5), Inches(2.8), Inches(1.0), NAVY2, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, bx + Inches(0.15), Inches(3.62), Inches(2.55), Inches(0.85), [[("✓ ", 13, TEAL, True), (c, 12.5, RGBColor(0xD7,0xE0,0xEC), False)]])
text(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.6),
     [[("Image légitime → ", 17, INK, False), ("✅ le pod tourne", 18, TEALD, True),
       ("   ·   image piégée → ", 17, INK, False), ("❌ bloquée à l'entrée", 18, RED, True)]])
# panorama outils
text(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(0.4), [("Panorama des outils vus :", 13, GREY, True)])
tools = [("Syft", TEALD), ("Grype", AMBER), ("cosign/Sigstore", NAVY2), ("Kyverno", RED)]
x = Inches(0.6)
for name, col in tools:
    chip(s, x, Inches(6.15), Inches(2.6), name, col, fs=12); x = x + Inches(2.75)
footer(s, 11)
notes(s, [
    "⏱️ 24:00–26:00 (2 min).",
    "Rejouer mentalement toute la chaîne d'un trait pour ancrer. Chaque étape produit une preuve ; "
    "l'admission les consomme.",
    "Si le temps le permet : lancer la démo complète (déployer image signée = OK, image non signée = refus).",
    "Rappeler les alternatives une dernière fois : Trivy (scan), Notation/Notary v2 (signature, écosystème "
    "Azure/ACR), OPA-Gatekeeper (admission). Le concept prime sur l'outil.",
])

# =================================================================== 12. À RETENIR / QUESTIONS
s = slide(); bg(s, NAVY)
rect(s, 0, Inches(3.25), W, Inches(0.06), TEAL)
text(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.9), [("À retenir", 34, WHITE, True)])
bullets(s, Inches(0.95), Inches(1.9), Inches(11.4), Inches(1.4), [
    "La supply chain est LA cible : rendre chaque maillon prouvable.",
    "Preuves accrochées au DIGEST : SBOM, signature, provenance.",
], size=18, color=RGBColor(0xD7,0xE0,0xEC), gap=8)
text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.6),
     [[("On ne fait pas confiance — ", 24, WHITE, True), ("on vérifie.", 24, TEAL, True)]])
bullets(s, Inches(0.95), Inches(4.4), Inches(11.4), Inches(1.8), [
    "Syft/Grype (connaître & auditer) · cosign/Sigstore (prouver l'origine) · Kyverno (bloquer l'inconnu).",
    "SLSA = boussole de maturité ; visez L2, soyez honnêtes sur les limites.",
    "Enforce, pas Audit, quand vous voulez réellement protéger.",
], size=16, color=RGBColor(0xC7,0xD2,0xE0), gap=9)
text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5), [("Questions ?  →  place au TP.", 16, TEAL, True)])
footer(s, 12, dark=True)
notes(s, [
    "⏱️ 26:00–28:00 + Q/R jusqu'à 30:00.",
    "Résumer en 3 phrases. Reprendre le slogan « on ne fait pas confiance, on vérifie ».",
    "Ouvrir les questions. Questions probables : coût/perf de la vérif à l'admission (négligeable, cache) ; "
    "que faire des images tierces non signées (les signer soi-même après vérif, ou les mettre en exception "
    "tracée) ; keyless sans Internet (possible en self-hosted Sigstore).",
    "Enchaîner sur le TP : « vous allez construire exactement cette chaîne sur votre fork ».",
])

# ------------------------------------------------------------------- save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "cours-supply-chain-security.pptx")
prs.save(out)
print("OK ->", out, f"({len(prs.slides._sldIdLst)} slides, avec notes présentateur)")
